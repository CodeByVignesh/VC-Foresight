from __future__ import annotations

from io import BytesIO
from pathlib import Path

from fastapi import UploadFile
from pypdf import PdfReader
from pptx import Presentation

from app.models import DocumentContent
from app.utils.text import normalize_whitespace, truncate_text


class DocumentParser:
    async def parse(self, uploaded_file: UploadFile | None) -> DocumentContent:
        if uploaded_file is None:
            return DocumentContent()

        filename = uploaded_file.filename or "uploaded-document"
        suffix = Path(filename).suffix.lower()
        raw_bytes = await uploaded_file.read()
        await uploaded_file.close()

        if not raw_bytes:
            return DocumentContent(
                filename=filename,
                limitations=[f"Uploaded document {filename} was empty."],
            )

        try:
            if suffix == ".pdf":
                text = self._extract_pdf_text(raw_bytes)
                document_type = "pdf"
            elif suffix == ".pptx":
                text = self._extract_pptx_text(raw_bytes)
                document_type = "pptx"
            else:
                return DocumentContent(
                    filename=filename,
                    limitations=[
                        f"Unsupported document type for {filename}. Only .pdf and .pptx are accepted."
                    ],
                )
        except Exception as exc:
            return DocumentContent(
                filename=filename,
                limitations=[f"Document parsing failed for {filename}: {exc}"],
            )

        cleaned_text = truncate_text(normalize_whitespace(text), 18_000)
        if not cleaned_text:
            return DocumentContent(
                filename=filename,
                document_type=document_type,
                limitations=[f"No readable text could be extracted from {filename}."],
            )

        return DocumentContent(
            filename=filename,
            document_type=document_type,
            text=cleaned_text,
            extracted=True,
        )

    def _extract_pdf_text(self, raw_bytes: bytes) -> str:
        reader = PdfReader(BytesIO(raw_bytes))
        return " ".join((page.extract_text() or "") for page in reader.pages)

    def _extract_pptx_text(self, raw_bytes: bytes) -> str:
        presentation = Presentation(BytesIO(raw_bytes))
        collected: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    collected.append(text)
        return " ".join(collected)
