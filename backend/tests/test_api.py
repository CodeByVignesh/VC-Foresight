from io import BytesIO
from pathlib import Path
import uuid

from fastapi.testclient import TestClient
from pptx import Presentation

from app.main import app
from app.models import (
    CompetitorMetrics,
    PatentMetrics,
    PortfolioCompanyCreate,
    ResearchMetrics,
    StartupSignals,
    WebsiteContent,
)
from app.scoring import CURRENT_YEAR
from app.services.crm_repository import CRMRepository
from app.services.llm_extractor import LLMExtractor
from app.services.openalex_client import OpenAlexClient
from app.services.patents_provider import PlaceholderPatentsProvider
from app.services.portfolio_repository import PortfolioRepository
from app.services.search_provider import PlaceholderSearchProvider
from app.services.website_fetcher import WebsiteFetcher


def test_root_and_health_routes() -> None:
    with TestClient(app) as client:
        root_response = client.get("/")
        assert root_response.status_code == 200
        root_payload = root_response.json()
        assert root_payload["service"] == "Startup Novelty API"
        assert root_payload["docs_url"] == "/docs"

        health_response = client.get("/health")
        assert health_response.status_code == 200
        assert health_response.json()["status"] == "ok"


def test_score_startup_accepts_required_document_with_optional_website(monkeypatch) -> None:
    async def fake_fetch_website(self, website: str) -> WebsiteContent:
        return WebsiteContent(
            url=website,
            title="Example AI",
            meta_description="AI platform for hospitals",
            text="Hospital automation workflow platform",
            fetched=True,
        )

    async def fake_fetch_research(self, sector: str, description: str) -> ResearchMetrics:
        return ResearchMetrics(
            recent_paper_count=3,
            publication_years=[CURRENT_YEAR, CURRENT_YEAR - 1, CURRENT_YEAR - 2],
            top_titles=["Hospital Workflow AI"],
            topics=["Health AI", "Clinical operations"],
            trend_ratio=0.5,
            latest_publication_year=CURRENT_YEAR,
            provider_confidence=0.8,
        )

    async def fake_search_patents(self, startup) -> PatentMetrics:
        return PatentMetrics(provider_confidence=0.0)

    async def fake_search_competitors(self, startup) -> CompetitorMetrics:
        return CompetitorMetrics(provider_confidence=0.0)

    async def fake_extract_signals(self, startup, website, document, research, patents, competitors):
        return (
            StartupSignals(
                product_category="HealthTech AI",
                target_customer="Hospitals",
                claimed_innovation="AI workflow automation for hospital operations",
                technical_keywords=["hospital", "automation", "workflow"],
                market_trends=["Operational efficiency", "Clinical AI"],
                evidence_summary=["Pitch deck and website both describe hospital workflow automation."],
                website_signal_present=True,
            ),
            [],
        )

    monkeypatch.setattr(WebsiteFetcher, "fetch", fake_fetch_website)
    monkeypatch.setattr(OpenAlexClient, "fetch_research", fake_fetch_research)
    monkeypatch.setattr(PlaceholderPatentsProvider, "search", fake_search_patents)
    monkeypatch.setattr(PlaceholderSearchProvider, "search", fake_search_competitors)
    monkeypatch.setattr(LLMExtractor, "extract_signals", fake_extract_signals)

    db_path = Path(f"data/test_vc_portfolio_{uuid.uuid4().hex}.db")
    if db_path.exists():
        db_path.unlink()
    repository = PortfolioRepository(db_path)
    repository.init_db()
    repository.add_company(
        PortfolioCompanyCreate(
            company_name="Hospital Flow",
            website="https://hospitalflow.ai",
            sector="HealthTech AI",
            thesis="Workflow automation software for hospitals and clinical operations",
            keywords=["hospital", "workflow", "automation", "clinical"],
        )
    )
    crm_db_path = Path(f"data/test_vc_crm_{uuid.uuid4().hex}.db")
    if crm_db_path.exists():
        crm_db_path.unlink()
    crm_repository = CRMRepository(crm_db_path)
    crm_repository.init_db()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Example AI"
    slide.placeholders[1].text = "AI platform for hospital workflow automation"
    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)

    with TestClient(app) as client:
        app.state.portfolio_repository = repository
        app.state.crm_repository = crm_repository
        response = client.post(
            "/score-startup",
            data={
                "sector": "HealthTech AI",
                "meeting_notes": "Founder emphasized workflow automation and hospital deployment plans.",
                "pitch_date": "2026-05-01",
                "deal_status": "screening",
                "funding_status": "seeking",
                "founder_names": "Alice Doe,Bob Roe",
            },
            files={
                "supporting_document": (
                    "deck.pptx",
                    buffer.getvalue(),
                    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )
            },
        )

    assert response.status_code == 200
    payload = response.json()
    assert payload["startup_name"] == "Deck"
    assert payload["predicted_domain"] == "HealthTech AI"
    assert payload["novelty_score"] >= 0
    assert 0 <= payload["novelty_score_10"] <= 10
    assert 0 <= payload["fit_score_10"] <= 10
    assert 0 <= payload["foresight_score_10"] <= 10
    assert payload["crm_record"]["recorded"] is True
    assert payload["portfolio_check"]["checked"] is True
    assert payload["portfolio_check"]["has_similar_investment"] is True
    assert payload["portfolio_check"]["top_matches"][0]["company_name"] == "Hospital Flow"
    assert any(item["source"] == "document_pptx" for item in payload["evidence"])
    assert any(item["source"] == "meeting_notes" for item in payload["evidence"])
    crm_companies = crm_repository.list_companies()
    crm_pitches = crm_repository.list_pitches()
    assert len(crm_companies) == 1
    assert len(crm_pitches) == 1
    assert crm_companies[0].predicted_domain == "HealthTech AI"
    assert crm_pitches[0].predicted_domain == "HealthTech AI"
    assert crm_pitches[0].deal_status == "screening"

    if db_path.exists():
        db_path.unlink()
    if crm_db_path.exists():
        crm_db_path.unlink()


def test_score_startup_requires_supporting_document() -> None:
    with TestClient(app) as client:
        response = client.post(
            "/score-startup",
            data={"startup_name": "No Deck Startup"},
        )

    assert response.status_code == 422


def test_portfolio_company_endpoints() -> None:
    db_path = Path(f"data/test_vc_portfolio_endpoints_{uuid.uuid4().hex}.db")
    if db_path.exists():
        db_path.unlink()
    repository = PortfolioRepository(db_path)
    repository.init_db()

    with TestClient(app) as client:
        app.state.portfolio_repository = repository
        create_response = client.post(
            "/portfolio-companies",
            json={
                "company_name": "Infra Brain",
                "website": "https://infra.example",
                "sector": "Developer Tools",
                "country": "Germany",
                "thesis": "AI infra for engineering teams",
                "keywords": ["developer", "infrastructure", "ai"],
            },
        )
        assert create_response.status_code == 200

        list_response = client.get("/portfolio-companies")
        assert list_response.status_code == 200
        companies = list_response.json()
        assert len(companies) == 1
        assert companies[0]["company_name"] == "Infra Brain"

    if db_path.exists():
        db_path.unlink()


def test_crm_endpoints_and_summary() -> None:
    crm_db_path = Path(f"data/test_vc_crm_endpoints_{uuid.uuid4().hex}.db")
    if crm_db_path.exists():
        crm_db_path.unlink()
    crm_repository = CRMRepository(crm_db_path)
    crm_repository.init_db()

    with TestClient(app) as client:
        app.state.crm_repository = crm_repository
        company_response = client.post(
            "/crm/companies",
            json={
                "company_name": "Infra Brain",
                "website": "https://infra.example",
                "sector": "Developer Tools",
                "country": "Germany",
                "description": "AI infra for engineering teams",
                "founder_names": ["Alice"],
                "keywords": ["developer", "infrastructure", "ai"],
            },
        )
        assert company_response.status_code == 200
        company_id = company_response.json()["id"]

        pitch_response = client.post(
            "/crm/pitches",
            json={
                "company_id": company_id,
                "pitch_date": "2026-05-01",
                "deal_status": "partner_review",
                "funding_status": "in_discussion",
                "round_name": "Seed",
                "amount_requested_usd": 750000,
                "source": "warm_intro",
                "notes": "Strong technical team",
            },
        )
        assert pitch_response.status_code == 200

        pitches_response = client.get("/crm/pitches")
        assert pitches_response.status_code == 200
        assert len(pitches_response.json()) == 1

        summary_response = client.get("/crm/summary")
        assert summary_response.status_code == 200
        summary = summary_response.json()
        assert summary["total_companies"] == 1
        assert summary["total_pitches"] == 1
        assert summary["domain_counts"][0]["label"] == "Developer Tools AI"
        assert summary["deal_status_counts"][0]["label"] == "partner_review"

    if crm_db_path.exists():
        crm_db_path.unlink()
